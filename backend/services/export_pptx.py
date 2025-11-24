from pptx import Presentation

def export_pptx(project, sections):
    prs = Presentation()

    for sec in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = sec.title
        slide.placeholders[1].text = sec.content or ""

    path = f"exports/{project.title.replace(' ', '_')}.pptx"
    prs.save(path)
    return path
